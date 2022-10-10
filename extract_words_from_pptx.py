# this utils extract words from pptx
import os
import glob
from pptx import Presentation




def get_word_from_ppt_folder(pptx_folder_path):
    if not os.path.exists(pptx_folder_path):
        raise Exception(
            ('current folder: ' + pptx_folder_path + ' does not exists. Check the spelling of your specified folder'))


    # TODO: check ppt extension
    extension_name_list = ['.pptx']
    filenames_with_extensions = []
    for extension_name in extension_name_list:
        for filename in glob.iglob(pptx_folder_path + '**\*' + extension_name, recursive=True):
            # skip temp files
            filename_last = filename.split('\\')[-1]
            if "~$" not in filename_last:
                filenames_with_extensions.append(filename)

    print(filenames_with_extensions)
    if len(filenames_with_extensions) == 0:
        raise Exception(('current folder: ' + pptx_folder_path + ' does not contain specified extensions'))

    all_text_dict = {}
    for eachfile in filenames_with_extensions:
        print(eachfile)
        eachfile_txt = []
        prs = Presentation(eachfile)
        for slide_num, slide in enumerate(prs.slides):
            for shape_num, shape in enumerate(slide.shapes):
                if hasattr(shape, "text"):

                    eachfile_txt.append(shape.text)

                    """
         
                    if search_term in shape.text.lower():
                        print(('Slide filename: ' + str(eachfile) + ', slide number: ' + str(slide_num + 1) +
                               ', shape number: ' + str(shape_num + 1)))
                    """

        all_text_dict[eachfile] = eachfile_txt


    return all_text_dict



if __name__ == "__main__":

    test_pptx_path = r"C:\\Users\maomao\Downloads"

    text_dictionary_from_folder = get_word_from_ppt_folder(test_pptx_path)

    test_file = list(text_dictionary_from_folder.keys())[0]

    test_text = text_dictionary_from_folder[test_file]

    devide_str = ' '
    test_concatinate = devide_str.join((str(n) for n in test_text if len(n)>0))

    # minimal example
    from keybert import KeyBERT

    kw_model = KeyBERT()
    keywords = kw_model.extract_keywords(test_concatinate)
    print(keywords)

    # extract phrases
    keyphrases = kw_model.extract_keywords(test_concatinate, keyphrase_ngram_range=(1, 4), stop_words='english',
                              use_mmr=True, diversity=0.7)
    print(keyphrases)


