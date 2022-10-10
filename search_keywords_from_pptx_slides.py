# search key words in slides

import sys

# To install the package, use the following command:
# pip install python-pptx

from pptx import Presentation
import glob
import os




if __name__ == '__main__':

    #test_pptx_path = "Test Path with pptx slides"
    #search_term = 'results'
    test_pptx_path = sys.argv[1]
    search_term = sys.argv[2]


    if not os.path.exists(test_pptx_path):
        raise Exception(('current folder: '+test_pptx_path + ' does not exists. Check the spelling of your specified folder'))


    #TODO: check ppt extension
    extension_name_list = ['.pptx']    
    filenames_with_extensions= []
    for extension_name in extension_name_list:
        for filename in glob.iglob(test_pptx_path + '**\*'+extension_name, recursive=True):
            filenames_with_extensions.append(filename)

    if len(filenames_with_extensions) == 0:
        raise Exception(('current folder: '+test_pptx_path + ' does not contain specified extensions'))

    for eachfile in filenames_with_extensions:
        prs = Presentation(eachfile)
        for slide_num, slide in enumerate(prs.slides):
            for shape_num, shape in enumerate(slide.shapes):
                if hasattr(shape, "text"):

                    if search_term in shape.text.lower():
                        print(('Slide filename: '+str(eachfile)+', slide number: ' + str(slide_num+1) +
                               ', shape number: ' + str(shape_num+1)) )


